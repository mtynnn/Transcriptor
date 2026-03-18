import os
from pptx import Presentation
import yake
from pypdf import PdfReader

def extract_text_from_context(file_path: str) -> str:
    """Extrae todo el texto de un archivo PPTX o PDF."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Archivo no encontrado: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()
    text_runs = []

    if ext in ['.pptx', '.ppt']:
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
        except Exception as e:
            raise ValueError(f"No se pudo leer el archivo PPTX (puede estar corrupto): {e}")
    elif ext == '.pdf':
        try:
            reader = PdfReader(file_path)
            if reader.is_encrypted:
                raise ValueError("El PDF está encriptado con contraseña. No se puede extraer texto.")
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_runs.append(page_text)
        except ValueError:
            raise
        except Exception as e:
            raise ValueError(f"No se pudo leer el PDF (puede estar corrupto): {e}")

    return "\n".join(text_runs)


def extract_keywords(text: str, max_keywords: int = 40) -> list[str]:
    """
    Extrae las palabras clave más relevantes usando YAKE.
    Retorna una lista de palabras clave (strings).
    """
    if not text.strip():
        return []

    kw_extractor = yake.KeywordExtractor(
        lan="es",
        n=1,
        dedupLim=0.9,
        top=max_keywords,
        features=None
    )
    
    keywords_scored = kw_extractor.extract_keywords(text)
    keywords = [kw for kw, score in keywords_scored]
    return keywords

def get_whisper_initial_prompt(file_path: str, max_keywords: int = 40) -> str:
    """Extrae el texto del archivo, saca las palabras clave y las formatea como prompt."""
    try:
        text = extract_text_from_context(file_path)
        keywords = extract_keywords(text, max_keywords)
        prompt = ", ".join(keywords)
        return prompt
    except Exception as e:
        print(f"Error procesando archivo de contexto: {e}")
        return ""
